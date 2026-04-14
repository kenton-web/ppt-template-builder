#!/usr/bin/env python3
"""
build_deck.py — Branded PPTX deck builder.

Extracts visual chrome (header/footer bar images, backgrounds, fonts, colors)
from an existing .pptx template and builds fresh data-filled slides that
inherit the brand identity exactly.

Usage as library:
    from build_deck import DeckBuilder
    b = DeckBuilder("template.pptx", "output.pptx")
    b.add_cover(...)
    b.add_kpi(...)
    b.save()

Usage as CLI inspector:
    python3 build_deck.py inspect /path/to/template.pptx
"""

import sys
import os
import io
import json
import subprocess

# ── Auto-install deps ─────────────────────────────────────────────────────────
def _ensure_deps():
    for pkg, imp in [("python-pptx", "pptx"), ("pillow", "PIL"), ("lxml", "lxml")]:
        try:
            __import__(imp)
        except ImportError:
            print(f"Installing {pkg}...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", "--quiet", pkg])

_ensure_deps()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


# ─────────────────────────────────────────────────────────────────────────────
# TEMPLATE INSPECTOR
# ─────────────────────────────────────────────────────────────────────────────

class TemplateProfile:
    """
    Holds everything extracted from the template:
      - Slide dimensions
      - Chrome image blobs (header bar, footer bar, backgrounds) with positions
      - Fonts and colors detected from text shapes
      - Derived content-area bounds
    """

    def __init__(self, prs: Presentation):
        self.slide_w = prs.slide_width
        self.slide_h = prs.slide_height
        self.slide_w_in = round(self.slide_w / 914400, 4)
        self.slide_h_in = round(self.slide_h / 914400, 4)
        self.layouts = [l.name for l in prs.slide_layouts]

        # Chrome images keyed by (left_in, top_in) → {blob, ext, width_in, height_in}
        self._chrome: dict[tuple, dict] = {}

        # Fonts / colors
        self.font_heading = "Arial"
        self.font_body = "Arial"
        self.color_title_text = RGBColor(0x11, 0x18, 0x27)
        self.color_accent = RGBColor(0xEB, 0x41, 0x47)
        self.color_primary = RGBColor(0x1B, 0x3A, 0x6B)

        # Content area (derived after chrome extraction)
        self.content_top_in = 1.3
        self.content_bottom_in = None
        self.content_left_in = 1.0
        self.content_right_in = None
        self.content_width_in = None
        self.content_height_in = None

        # Title/cover layout type
        self.has_title_layout = False
        self.has_default_layout = False

        self._extract(prs)

    def _extract(self, prs: Presentation):
        TITLE_NAMES = {"title", "cover", "section"}
        DEFAULT_NAMES = {"default", "content", "blank", "layout"}

        for layout in prs.slide_layouts:
            n = layout.name.lower()
            if any(t in n for t in TITLE_NAMES):
                self.has_title_layout = True
            if any(d in n for d in DEFAULT_NAMES):
                self.has_default_layout = True

        # Scan all slides to collect chrome images and text styles
        header_bottoms = []
        footer_tops = []
        left_margins = []

        for slide in prs.slides:
            for shape in slide.shapes:
                l = shape.left / 914400
                t = shape.top / 914400
                w = shape.width / 914400
                h = shape.height / 914400

                if shape.shape_type == 13:  # PICTURE
                    key = (round(l, 2), round(t, 2))
                    if key not in self._chrome:
                        self._chrome[key] = {
                            "blob": shape.image.blob,
                            "ext":  shape.image.ext,
                            "left_in":   l,
                            "top_in":    t,
                            "width_in":  w,
                            "height_in": h,
                        }

                    # Header bar heuristic: near top, spans most of width
                    if t < 2.0 and w > self.slide_w_in * 0.5 and h < 2.0 and t > 0.1:
                        header_bottoms.append(t + h)
                        left_margins.append(l)
                    # Footer bar heuristic: near bottom, spans most of width
                    if t > self.slide_h_in * 0.7 and w > self.slide_w_in * 0.5 and h < 2.0:
                        footer_tops.append(t)

                # Extract text styles
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                color = run.font.color.rgb
                                if run.font.bold and run.font.size and run.font.size > Pt(20):
                                    self.color_title_text = color
                                if run.font.name and run.font.name not in ("", None):
                                    if run.font.bold:
                                        self.font_heading = run.font.name
                                    else:
                                        self.font_body = run.font.name
                            except Exception:
                                pass

        # Derive content area
        if header_bottoms:
            self.content_top_in = round(max(header_bottoms) + 0.15, 3)
        if footer_tops:
            self.content_bottom_in = round(min(footer_tops) - 0.15, 3)
        if left_margins:
            self.content_left_in = round(min(left_margins), 3)

        self.content_right_in = round(self.slide_w_in - self.content_left_in, 3)
        self.content_width_in = round(self.content_right_in - self.content_left_in, 3)
        if self.content_bottom_in:
            self.content_height_in = round(self.content_bottom_in - self.content_top_in, 3)

    def get_chrome(self, left_in: float, top_in: float) -> dict | None:
        """Look up a chrome image by approximate position."""
        for (l, t), img in self._chrome.items():
            if abs(l - left_in) < 0.1 and abs(t - top_in) < 0.1:
                return img
        return None

    def chrome_images(self) -> list[dict]:
        return list(self._chrome.values())

    def header_images(self) -> list[dict]:
        """Images near the top of the slide (header bars)."""
        return [img for img in self._chrome.values()
                if 0.1 < img["top_in"] < 2.0
                and img["width_in"] > self.slide_w_in * 0.5
                and img["height_in"] < 2.0]

    def footer_images(self) -> list[dict]:
        """Images near the bottom of the slide (footer bars)."""
        return [img for img in self._chrome.values()
                if img["top_in"] > self.slide_h_in * 0.7
                and img["width_in"] > self.slide_w_in * 0.5
                and img["height_in"] < 2.0]

    def background_images(self) -> list[dict]:
        """Full-bleed background images."""
        return [img for img in self._chrome.values()
                if img["left_in"] < 0.1 and img["top_in"] <= 0.05
                and img["width_in"] >= self.slide_w_in * 0.9]

    def to_dict(self) -> dict:
        chrome_summary = []
        for img in self._chrome.values():
            chrome_summary.append({
                "left_in": round(img["left_in"], 3),
                "top_in":  round(img["top_in"], 3),
                "width_in": round(img["width_in"], 3),
                "height_in": round(img["height_in"], 3),
                "size_bytes": len(img["blob"]),
                "ext": img["ext"],
                "role": self._classify_chrome(img),
            })
        return {
            "slide_w_in": self.slide_w_in,
            "slide_h_in": self.slide_h_in,
            "layouts": self.layouts,
            "fonts": {"heading": self.font_heading, "body": self.font_body},
            "colors": {
                "title_text": str(self.color_title_text),
                "accent": str(self.color_accent),
                "primary": str(self.color_primary),
            },
            "content_area": {
                "top":    self.content_top_in,
                "bottom": self.content_bottom_in,
                "left":   self.content_left_in,
                "right":  self.content_right_in,
                "width":  self.content_width_in,
                "height": self.content_height_in,
            },
            "chrome_images": chrome_summary,
        }

    def _classify_chrome(self, img: dict) -> str:
        if img["left_in"] < 0.1 and img["top_in"] <= 0.05:
            return "background"
        if img["top_in"] < 2.0 and img["width_in"] > self.slide_w_in * 0.5:
            return "header"
        if img["top_in"] > self.slide_h_in * 0.7:
            return "footer"
        return "decoration"


def inspect_template(path: str) -> TemplateProfile:
    """Load a template and return its TemplateProfile."""
    prs = Presentation(os.path.expanduser(path))
    return TemplateProfile(prs)


# ─────────────────────────────────────────────────────────────────────────────
# DECK BUILDER
# ─────────────────────────────────────────────────────────────────────────────

# Metric types that are "cost" (positive delta = bad = RED)
COST_METRICS = {"acos", "tacos", "cpc", "cpm", "spend", "cost", "fees", "stranded", "returns"}

TAG_COLORS = {
    "WIN":    (RGBColor(0xDC, 0xFC, 0xE7), RGBColor(0x16, 0xA3, 0x4A)),
    "RISK":   (RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0xDC, 0x26, 0x26)),
    "ACTION": (RGBColor(0xFF, 0xED, 0xD5), RGBColor(0xC2, 0x41, 0x0C)),
    "INFO":   (RGBColor(0xCC, 0xFB, 0xF1), RGBColor(0x0F, 0x76, 0x6E)),
}


class DeckBuilder:
    """
    Build a branded .pptx deck from a template.

    Extracts the template's visual chrome on init, strips the existing slides,
    then adds new slides via add_cover / add_section / add_kpi / add_chart /
    add_table / add_insights.  Call save() to write the file.
    """

    def __init__(self, template_path: str, output_path: str):
        self.template_path = os.path.expanduser(template_path)
        self.output_path   = os.path.expanduser(output_path)

        self.prs = Presentation(self.template_path)
        self.profile = TemplateProfile(self.prs)

        # Strip existing slides (drop_rel is required — just removing from
        # _sldIdLst leaves orphan parts that cause duplicate-name warnings)
        for sldId in list(self.prs.slides._sldIdLst):
            rId = sldId.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            )
            if rId:
                try:
                    self.prs.part.drop_rel(rId)
                except Exception:
                    pass
            self.prs.slides._sldIdLst.remove(sldId)

        assert len(self.prs.slides) == 0, "Slide strip failed"

        # Convenience aliases from profile
        p = self.profile
        self.W  = p.slide_w
        self.H  = p.slide_h
        self.LM = Inches(p.content_left_in)
        self.CW = Inches(p.content_width_in)
        self.CT = Inches(p.content_top_in)
        self.CB = Inches(p.content_bottom_in) if p.content_bottom_in else self.H - Inches(1.0)
        self.CH = self.CB - self.CT
        self.RM = self.LM + self.CW
        self.FH = self.profile.font_heading
        self.FB = self.profile.font_body
        self.C_TITLE    = p.color_title_text
        self.C_ACCENT   = p.color_accent
        self.C_PRIMARY  = p.color_primary
        self.C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
        self.C_LGRAY    = RGBColor(0xF3, 0xF4, 0xF6)
        self.C_MGRAY    = RGBColor(0x6B, 0x72, 0x80)
        self.C_ALTROW   = RGBColor(0xE8, 0xEA, 0xED)
        self.C_GREEN    = RGBColor(0x16, 0xA3, 0x4A)
        self.C_RED      = RGBColor(0xDC, 0x26, 0x26)
        self.C_ORANGE   = RGBColor(0xF9, 0x73, 0x16)
        self.C_CYAN     = RGBColor(0x06, 0xB6, 0xD4)
        self.C_BORDER   = RGBColor(0xD1, 0xD5, 0xDB)
        self.SERIES_COLORS = [self.C_PRIMARY, self.C_ORANGE, self.C_CYAN, self.C_MGRAY]

        self.slide_count = 0

    # ── Layout helper ─────────────────────────────────────────────────────────
    def _get_layout(self, name_hint: str):
        for layout in self.prs.slide_layouts:
            if name_hint.lower() in layout.name.lower():
                return layout
        return self.prs.slide_layouts[min(1, len(self.prs.slide_layouts) - 1)]

    # ── Drawing primitives ────────────────────────────────────────────────────
    def _add_img(self, slide, blob, ext, left, top, width, height):
        return slide.shapes.add_picture(io.BytesIO(blob), left, top, width, height)

    def _add_txt(self, slide, text, left, top, width, height,
                 size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font=None, italic=False, wrap=True):
        font = font or self.FB
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        if color:
            r.font.color.rgb = color
        return tb

    def _add_rect(self, slide, left, top, width, height,
                  fill=None, line_color=None, line_width=Pt(0)):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = line_width
        else:
            shape.line.fill.background()
        return shape

    # ── Chrome appliers ───────────────────────────────────────────────────────
    def _chrome_default(self, slide, title_text: str):
        """Apply header bar + footer bar + title text to a DEFAULT-type slide."""
        for img in self.profile.header_images():
            self._add_img(
                slide, img["blob"], img["ext"],
                Inches(img["left_in"]), Inches(img["top_in"]),
                Inches(img["width_in"]), Inches(img["height_in"]),
            )
        for img in self.profile.footer_images():
            self._add_img(
                slide, img["blob"], img["ext"],
                Inches(img["left_in"]), Inches(img["top_in"]),
                Inches(img["width_in"]), Inches(img["height_in"]),
            )
        # Title text overlaid on header
        self._add_txt(
            slide, title_text,
            self.LM - Inches(0.01), Inches(self.profile.content_top_in - 0.72),
            self.CW * 0.8, Inches(0.65),
            size=max(18, int(self.profile.slide_h_in * 2.8)),
            bold=True, color=self.C_TITLE, font=self.FH,
        )

    def _chrome_title(self, slide, use_alt_bg=False):
        """Apply full-bleed background + logo + decorations to a TITLE-type slide."""
        bgs = self.profile.background_images()
        if not bgs:
            # Fallback: solid navy background
            self._add_rect(slide, 0, 0, self.W, self.H, fill=self.C_PRIMARY)
            return

        # If multiple background images, try to alternate (alt for section dividers)
        bg = bgs[min(1 if use_alt_bg and len(bgs) > 1 else 0, len(bgs) - 1)]
        self._add_img(slide, bg["blob"], bg["ext"],
                      Inches(0), Inches(0), self.W, self.H)

        # Non-background, non-header, non-footer decorations (logo, deco images)
        for img in self.profile.chrome_images():
            role = self.profile._classify_chrome(img)
            if role == "decoration":
                self._add_img(
                    slide, img["blob"], img["ext"],
                    Inches(img["left_in"]), Inches(img["top_in"]),
                    Inches(img["width_in"]), Inches(img["height_in"]),
                )

    # ── Metric formatting ─────────────────────────────────────────────────────
    @staticmethod
    def _fmt(value, mtype: str) -> str:
        try:
            v = float(value)
        except (TypeError, ValueError):
            return str(value)
        if mtype == "currency":
            if v >= 1_000_000:   return f"${v/1_000_000:.1f}M"
            elif v >= 1_000:     return f"${v:,.0f}"
            else:                return f"${v:.2f}"
        elif mtype == "multiplier":  return f"{v:.2f}x"
        elif mtype in ("percent", "ratio"): return f"{v:.1f}%"
        elif mtype == "number":
            if v >= 1_000_000:   return f"{v/1_000_000:.1f}M"
            elif v >= 1_000:     return f"{v:,.0f}"
            return str(int(v))
        elif mtype == "days":    return f"{v:.0f}d"
        elif mtype == "rank":    return f"#{int(v)}"
        return str(value)

    def _delta_color(self, pct: float, name: str) -> RGBColor:
        is_cost = any(c in name.lower() for c in COST_METRICS)
        is_rank = "rank" in name.lower() or "bsr" in name.lower()
        if is_rank or is_cost:
            return self.C_GREEN if pct < 0 else self.C_RED
        return self.C_GREEN if pct > 0 else self.C_RED

    # ─────────────────────────────────────────────────────────────────────────
    # PUBLIC SLIDE BUILDERS
    # ─────────────────────────────────────────────────────────────────────────

    def add_cover(self, title: str, client_name: str, period: str,
                  report_type: str = "", presenter: str = "", date: str = ""):
        """Cover / title slide."""
        slide = self.prs.slides.add_slide(self._get_layout("TITLE"))
        self._chrome_title(slide, use_alt_bg=False)

        sw = self.profile.slide_w_in
        sh = self.profile.slide_h_in

        self._add_txt(slide, title,
                      Inches(sw * 0.04), Inches(sh * 0.36),
                      Inches(sw * 0.72), Inches(sh * 0.14),
                      size=int(sh * 5.0), bold=True,
                      color=self.C_WHITE, font=self.FH)

        self._add_txt(slide, client_name,
                      Inches(sw * 0.04), Inches(sh * 0.50),
                      Inches(sw * 0.72), Inches(sh * 0.10),
                      size=int(sh * 3.2), bold=False,
                      color=self.C_WHITE, font=self.FH)

        self._add_txt(slide, period,
                      Inches(sw * 0.04), Inches(sh * 0.60),
                      Inches(sw * 0.72), Inches(sh * 0.07),
                      size=int(sh * 2.3), bold=False,
                      color=RGBColor(0xCB, 0xD5, 0xE1), font=self.FB)

        if report_type:
            self._add_txt(slide, report_type,
                          Inches(sw * 0.04), Inches(sh * 0.67),
                          Inches(sw * 0.72), Inches(sh * 0.06),
                          size=int(sh * 1.8), bold=False,
                          color=RGBColor(0x9C, 0xA3, 0xAF), font=self.FB)

        if presenter or date:
            byline = f"{presenter}  ·  {date}".strip(" ·")
            self._add_txt(slide, byline,
                          Inches(sw * 0.04), Inches(sh * 0.87),
                          Inches(sw * 0.6), Inches(sh * 0.05),
                          size=int(sh * 1.4), bold=False,
                          color=RGBColor(0x9C, 0xA3, 0xAF), font=self.FB)

        self.slide_count += 1
        print(f"  ✦ Cover: {title} — {client_name}")

    def add_section(self, section_title: str, subtitle: str = "", number: str = ""):
        """Section divider slide (uses TITLE layout / alt background)."""
        slide = self.prs.slides.add_slide(self._get_layout("TITLE"))
        self._chrome_title(slide, use_alt_bg=True)

        sw = self.profile.slide_w_in
        sh = self.profile.slide_h_in

        if number:
            self._add_txt(slide, number,
                          Inches(sw * 0.04), Inches(sh * 0.32),
                          Inches(sw * 0.15), Inches(sh * 0.12),
                          size=int(sh * 6.5), bold=True,
                          color=self.C_ACCENT, font=self.FH)

        self._add_txt(slide, section_title,
                      Inches(sw * 0.04), Inches(sh * 0.44),
                      Inches(sw * 0.85), Inches(sh * 0.14),
                      size=int(sh * 4.6), bold=True,
                      color=self.C_WHITE, font=self.FH)

        if subtitle:
            self._add_txt(slide, subtitle,
                          Inches(sw * 0.04), Inches(sh * 0.58),
                          Inches(sw * 0.75), Inches(sh * 0.07),
                          size=int(sh * 2.1), bold=False,
                          color=RGBColor(0xCB, 0xD5, 0xE1), font=self.FB)

        self.slide_count += 1
        print(f"  ✦ Section: {number} — {section_title}")

    def add_kpi(self, slide_title: str, metrics: list[dict], period_label: str = ""):
        """
        KPI scorecard grid.
        Each metric dict: {name, value, metric_type, delta_pct (optional)}
        """
        slide = self.prs.slides.add_slide(self._get_layout("DEFAULT"))
        self._chrome_default(slide, slide_title)

        if period_label:
            self._add_txt(slide, period_label,
                          self.LM, self.CT + Inches(0.0),
                          self.CW, Inches(0.3),
                          size=11, bold=False, color=self.C_MGRAY)

        n = len(metrics)
        cols = min(4, n)
        rows = (n + cols - 1) // cols
        gap  = Inches(0.14)
        y0   = self.CT + Inches(0.35)
        card_w = (self.CW - gap * (cols - 1)) / cols
        card_h = (self.CH - Inches(0.35) - gap * (rows - 1)) / rows

        for i, m in enumerate(metrics):
            col = i % cols
            row = i // cols
            x = self.LM + col * (card_w + gap)
            y = y0 + row * (card_h + gap)

            # Card background
            self._add_rect(slide, x, y, card_w, card_h, fill=self.C_LGRAY)
            # Accent top bar
            self._add_rect(slide, x, y, card_w, Inches(0.055), fill=self.C_PRIMARY)

            ip = Inches(0.14)
            name_h = card_h * 0.28
            val_h  = card_h * 0.42
            dlt_h  = card_h * 0.28

            self._add_txt(slide, m["name"],
                          x + ip, y + ip * 0.6, card_w - ip * 2, name_h,
                          size=max(9, int(self.profile.slide_h_in * 1.15)),
                          bold=False, color=self.C_MGRAY)

            val_str = self._fmt(m["value"], m.get("metric_type", "number"))
            self._add_txt(slide, val_str,
                          x + ip, y + name_h, card_w - ip * 2, val_h,
                          size=max(16, int(self.profile.slide_h_in * 2.5)),
                          bold=True, color=self.C_TITLE, font=self.FH)

            if "delta_pct" in m and m["delta_pct"] is not None:
                d = float(m["delta_pct"])
                dc = self._delta_color(d, m["name"])
                sign = "▲" if d >= 0 else "▼"
                self._add_txt(slide, f"{sign} {abs(d):.1f}%  vs prior",
                              x + ip, y + card_h - dlt_h, card_w - ip * 2, dlt_h,
                              size=max(8, int(self.profile.slide_h_in * 0.95)),
                              bold=True, color=dc)

        self.slide_count += 1
        print(f"  ✦ KPI: {slide_title} ({n} metrics, {rows}×{cols})")

    def add_chart(self, slide_title: str, chart_data: dict,
                  chart_type: str = "column", subtitle: str = ""):
        """
        Trend chart.
        chart_data: {categories: [...], series: [{name, values}, ...]}
        chart_type: column | bar | line | bar_stacked
        """
        slide = self.prs.slides.add_slide(self._get_layout("DEFAULT"))
        self._chrome_default(slide, slide_title)

        if subtitle:
            self._add_txt(slide, subtitle,
                          self.LM, self.CT, self.CW, Inches(0.3),
                          size=11, bold=False, color=self.C_MGRAY)

        ct_map = {
            "column":      XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar":         XL_CHART_TYPE.BAR_CLUSTERED,
            "line":        XL_CHART_TYPE.LINE,
            "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
        }
        ct = ct_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        cd = ChartData()
        cd.categories = chart_data["categories"]
        for s in chart_data["series"]:
            cd.add_series(s["name"], s["values"])

        chart_top = self.CT + (Inches(0.4) if subtitle else Inches(0.1))
        chart_h   = self.CB - chart_top

        chart_obj = slide.shapes.add_chart(
            ct, self.LM, chart_top, self.CW, chart_h, cd
        ).chart

        chart_obj.has_legend = True
        chart_obj.legend.position = 2  # bottom
        chart_obj.legend.include_in_layout = False

        for i, series in enumerate(chart_obj.series):
            if i < len(self.SERIES_COLORS):
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = self.SERIES_COLORS[i]

        self.slide_count += 1
        print(f"  ✦ Chart: {slide_title} ({len(chart_data['series'])} series, {chart_type})")

    def add_table(self, slide_title: str, column_headers: list[str],
                  rows: list[list], subtitle: str = "",
                  highlight_col: int | None = None, sort_note: str = ""):
        """
        Data table with navy header row, alternate row shading.
        """
        slide = self.prs.slides.add_slide(self._get_layout("DEFAULT"))
        self._chrome_default(slide, slide_title)

        if subtitle:
            self._add_txt(slide, subtitle,
                          self.LM, self.CT, self.CW, Inches(0.3),
                          size=11, bold=False, color=self.C_MGRAY)

        n_rows = len(rows) + 1
        n_cols = len(column_headers)
        tbl_top = self.CT + (Inches(0.4) if subtitle else Inches(0.15))
        tbl_h   = self.CB - tbl_top - (Inches(0.3) if sort_note else Inches(0.05))

        tbl = slide.shapes.add_table(n_rows, n_cols, self.LM, tbl_top, self.CW, tbl_h).table

        # Distribute column widths: first col wider, rest equal
        first_w = int(self.CW * 0.32)
        second_w = int(self.CW * 0.08) if n_cols > 2 else 0
        rest_w = int((self.CW - first_w - second_w) / max(1, n_cols - (2 if second_w else 1)))
        for i in range(n_cols):
            if i == 0:
                tbl.columns[i].width = first_w
            elif i == 1 and second_w:
                tbl.columns[i].width = second_w
            else:
                tbl.columns[i].width = rest_w

        C_HL = RGBColor(0xEB, 0xF3, 0xFF)

        # Header row
        for j, hdr in enumerate(column_headers):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.C_PRIMARY
            tf = cell.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = hdr
            r.font.size = Pt(max(9, int(self.profile.slide_h_in * 1.05)))
            r.font.bold = True
            r.font.color.rgb = self.C_WHITE
            r.font.name = self.FH

        # Data rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i + 1, j)
                cell.fill.solid()
                if j == highlight_col:
                    cell.fill.fore_color.rgb = C_HL
                elif i % 2 == 1:
                    cell.fill.fore_color.rgb = self.C_ALTROW
                else:
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                tf = cell.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                r = p.add_run()
                r.text = str(val)
                r.font.size = Pt(max(8, int(self.profile.slide_h_in * 0.98)))
                r.font.bold = (j == highlight_col)
                r.font.color.rgb = self.C_TITLE
                r.font.name = self.FB

        if sort_note:
            self._add_txt(slide, sort_note,
                          self.LM, self.CB - Inches(0.28), self.CW, Inches(0.25),
                          size=9, bold=False, color=self.C_MGRAY,
                          align=PP_ALIGN.RIGHT)

        self.slide_count += 1
        print(f"  ✦ Table: {slide_title} ({len(rows)} rows × {n_cols} cols)")

    def add_insights(self, slide_title: str, bullets: list[dict],
                     highlight_metric: dict | None = None):
        """
        Insights slide with WIN / RISK / ACTION / INFO tag pills.
        Each bullet: {lead, body, tag}
        Optional highlight_metric: {label, value, caption}
        """
        slide = self.prs.slides.add_slide(self._get_layout("DEFAULT"))
        self._chrome_default(slide, slide_title)

        n = len(bullets)
        has_callout = highlight_metric is not None
        bullet_w = self.CW * (0.70 if has_callout else 1.0)
        gap = Inches(0.14)
        bh  = (self.CH - gap * (n - 1)) / n
        y0  = self.CT + Inches(0.05)
        ip  = Inches(0.14)

        for i, b in enumerate(bullets):
            tag = b.get("tag", "INFO").upper()
            bg_c, border_c = TAG_COLORS.get(tag, TAG_COLORS["INFO"])
            bx = self.LM
            by = y0 + i * (bh + gap)

            self._add_rect(slide, bx, by, bullet_w, bh,
                           fill=bg_c, line_color=border_c, line_width=Pt(1))

            # Tag pill background
            pill_w = Inches(0.65)
            pill_h = Inches(0.26)
            self._add_rect(slide, bx + ip, by + ip * 0.7, pill_w, pill_h, fill=border_c)
            self._add_txt(slide, tag,
                          bx + ip, by + ip * 0.7, pill_w, pill_h,
                          size=max(7, int(self.profile.slide_h_in * 0.8)),
                          bold=True, color=self.C_WHITE,
                          font=self.FH, align=PP_ALIGN.CENTER)

            lead_top = by + ip * 0.7 + pill_h + ip * 0.4
            lead_h   = bh * 0.32
            body_top = lead_top + lead_h
            body_h   = by + bh - body_top - ip * 0.5

            self._add_txt(slide, b["lead"],
                          bx + ip, lead_top, bullet_w - ip * 2, lead_h,
                          size=max(11, int(self.profile.slide_h_in * 1.35)),
                          bold=True, color=self.C_TITLE, font=self.FH)

            self._add_txt(slide, b.get("body", ""),
                          bx + ip, body_top, bullet_w - ip * 2, body_h,
                          size=max(9, int(self.profile.slide_h_in * 1.05)),
                          bold=False, color=self.C_MGRAY, wrap=True)

        # Callout box
        if has_callout:
            cx = self.LM + bullet_w + Inches(0.15)
            cw = self.CW - bullet_w - Inches(0.15)
            self._add_rect(slide, cx, y0, cw, self.CH, fill=self.C_PRIMARY)
            label_size = max(10, int(self.profile.slide_h_in * 1.2))
            val_size   = max(24, int(self.profile.slide_h_in * 4.3))
            cap_size   = max(9,  int(self.profile.slide_h_in * 1.1))

            self._add_txt(slide, highlight_metric["label"],
                          cx + Inches(0.15), y0 + self.CH * 0.08, cw - Inches(0.3), Inches(0.5),
                          size=label_size, bold=False,
                          color=RGBColor(0x9C, 0xA3, 0xAF),
                          font=self.FB, align=PP_ALIGN.CENTER)
            self._add_txt(slide, highlight_metric["value"],
                          cx + Inches(0.1), y0 + self.CH * 0.22, cw - Inches(0.2), self.CH * 0.35,
                          size=val_size, bold=True, color=self.C_ORANGE,
                          font=self.FH, align=PP_ALIGN.CENTER)
            self._add_txt(slide, highlight_metric.get("caption", ""),
                          cx + Inches(0.15), y0 + self.CH * 0.58, cw - Inches(0.3), self.CH * 0.38,
                          size=cap_size, bold=False,
                          color=RGBColor(0xCB, 0xD5, 0xE1),
                          font=self.FB, align=PP_ALIGN.CENTER, wrap=True)

        self.slide_count += 1
        print(f"  ✦ Insights: {slide_title} ({n} bullets)")

    def save(self) -> str:
        """Write the deck to disk and return the absolute path."""
        out = os.path.abspath(self.output_path)
        self.prs.save(out)
        print(f"\n✅ Saved: {out}  ({self.slide_count} slides)")
        return out


# ─────────────────────────────────────────────────────────────────────────────
# CLI: inspect command
# ─────────────────────────────────────────────────────────────────────────────

def cli_inspect(path: str):
    profile = inspect_template(path)
    print(json.dumps(profile.to_dict(), indent=2, default=str))


if __name__ == "__main__":
    if len(sys.argv) >= 3 and sys.argv[1] == "inspect":
        cli_inspect(sys.argv[2])
    else:
        print("Usage: python3 build_deck.py inspect /path/to/template.pptx")
        sys.exit(1)
